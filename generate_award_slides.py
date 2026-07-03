#!/usr/bin/env python3
from __future__ import annotations

import copy
import re
import sys
from collections import defaultdict

import openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.parts.slide import SlidePart

HEADER_MAP = {
    "award category": "award_category",
    "nominee name": "nominee_name",
    "winner name": "winner_name",
    "zone": "zone",
    "placeholder x": "placeholder_x",
    "placeholder y": "placeholder_y",
    "placeholder z": "placeholder_z",
}

TOKENS = {
    "<<zone>>": "ZONE",
    "<<award category>>": "AWARD CATEGORY",
    "<<nominees>>": "NOMINEES",
    "<<winner>>": "WINNER",
    "<<nominee name>>": "NOMINEES",
    "<<winner name>>": "WINNER",
    "<<nominees-word>>": "NOMINEES_WORD",
    "<<winner-word>>": "WINNER_WORD",
    "<<placeholder x>>": "PLACEHOLDER_X",
    "<<placeholder y>>": "PLACEHOLDER_Y",
    "<<placeholder z>>": "PLACEHOLDER_Z",
}


def norm(v):
    return re.sub(r"\s+", " ", "" if v is None else str(v)).strip().lower()


def clean(v):
    return "" if v is None else re.sub(r"\s+", " ", str(v)).strip()


def read_rows(xlsx_path):
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    rows = []
    for ws in wb.worksheets:
        data = list(ws.iter_rows(values_only=True))
        if not data:
            continue
        headers = [clean(x) for x in data[0]]
        idx = {}
        for i, h in enumerate(headers):
            key = HEADER_MAP.get(norm(h))
            if key:
                idx[key] = i
        for r in data[1:]:
            if not any(v is not None and str(v).strip() for v in r):
                continue
            def g(key):
                i = idx.get(key)
                return clean(r[i]) if i is not None and i < len(r) else ""
            rows.append({
                "award_category": g("award_category"),
                "nominee_name": g("nominee_name"),
                "winner_name": g("winner_name"),
                "zone": g("zone") or ws.title,
                "placeholder_x": g("placeholder_x"),
                "placeholder_y": g("placeholder_y"),
                "placeholder_z": g("placeholder_z"),
            })
    return rows


def group_rows(rows):
    """Group rows by (zone, award_category), while also recording the order
    in which each category key FIRST appears in the original Excel sheet.
    This lets build_deck() emit slides in the same sequence as the Excel
    file, which is critical for the emcee's script order."""
    nominee_groups = defaultdict(list)
    winner_groups = defaultdict(list)
    key_order = []
    seen = set()
    for r in rows:
        key = (r["zone"], r["award_category"])
        if key not in seen:
            seen.add(key)
            key_order.append(key)
        if r["nominee_name"]:
            nominee_groups[key].append(r)
        if r["winner_name"]:
            winner_groups[key].append(r)
    return nominee_groups, winner_groups, key_order


def find_slide_index_with_token(prs, token):
    token = token.lower()
    for i, s in enumerate(prs.slides):
        txt = " ".join([shape.text for shape in s.shapes if hasattr(shape, "text") and shape.text]).lower()
        if token in txt:
            return i
    return None


def pick_template_indices(prs):
    nominee_idx = find_slide_index_with_token(prs, "<<nominee name>>")
    if nominee_idx is None:
        nominee_idx = find_slide_index_with_token(prs, "<<nominees>>")
    winner_idx = find_slide_index_with_token(prs, "<<winner name>>")
    if winner_idx is None:
        winner_idx = find_slide_index_with_token(prs, "<<winner>>")
    if nominee_idx is None and len(prs.slides) > 0:
        nominee_idx = 0
    if winner_idx is None and len(prs.slides) > 1:
        winner_idx = 1
    if winner_idx is None:
        winner_idx = nominee_idx
    if nominee_idx is None or winner_idx is None:
        raise RuntimeError("Template must contain at least one slide with nominee or winner placeholders.")
    return nominee_idx, winner_idx


def duplicate_slide(prs, index):
    """True slide clone at the OPC-part level: copies the slide XML part
    byte-for-byte (preserving every font, color, schemeClr reference,
    autofit setting, and layout/master linkage exactly), then registers
    it as a new slide part+relationship in the presentation. This is the
    only approach that guarantees 100% fidelity with the original
    template's theme/color-map resolution."""
    source_slide = prs.slides[index]
    source_part = source_slide.part

    new_slide_el = copy.deepcopy(source_part._element)
    partname = prs.part.package.next_partname("/ppt/slides/slide%d.xml")
    new_slide_part = SlidePart(partname, source_part.content_type, source_part.package, new_slide_el)

    # Copy every relationship using the SAME rId as the source part, so all
    # r:embed / r:link / r:id references inside the copied XML remain valid
    # without any remapping. We write directly into the internal _rels dict
    # using the same _Relationship value object python-pptx uses internally,
    # since the public API only supports auto-assigned rIds.
    from pptx.opc.package import _Relationship
    from pptx.opc.packuri import PackURI
    for old_rid, rel in source_part.rels.items():
        if rel.is_external:
            new_slide_part.rels._rels[old_rid] = _Relationship(
                new_slide_part.partname.baseURI,
                old_rid,
                rel.reltype,
                target_mode="External",
                target=rel.target_ref,
            )
        else:
            new_slide_part.rels._rels[old_rid] = _Relationship(
                new_slide_part.partname.baseURI,
                old_rid,
                rel.reltype,
                target_mode="Internal",
                target=rel.target_part,
            )

    new_rid = prs.part.relate_to(
        new_slide_part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    )

    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(s.get("id")) for s in sldIdLst]
    new_id = max(existing_ids) + 1 if existing_ids else 256
    new_sldId = sldIdLst.makeelement(qn("p:sldId"), {"id": str(new_id)})
    new_sldId.set(qn("r:id"), new_rid)
    sldIdLst.append(new_sldId)

    return prs.slides[len(prs.slides) - 1]


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def replace_token_in_run_list(paragraph, token, value):
    """Replace a token that appears (possibly split across multiple runs)
    within a paragraph, WITHOUT touching run formatting (rPr) of the run
    that carries the token. If the token spans multiple runs, only the
    first run's formatting is kept and the rest of the matching runs are
    removed. Multi-line values create sibling paragraphs cloned from this
    paragraph's pPr/formatting so every line matches the template style."""
    runs = paragraph.runs
    full_text = "".join(r.text for r in runs)
    low = full_text.lower()
    pos = low.find(token)
    if pos == -1:
        return False

    lines = value.split("\n") if value else [""]
    first_line, rest_lines = lines[0], lines[1:]

    cum = 0
    start_run = end_run = None
    start_off = end_off = 0
    for i, r in enumerate(runs):
        rlen = len(r.text)
        if start_run is None and cum + rlen > pos:
            start_run = i
            start_off = pos - cum
        if cum + rlen >= pos + len(token):
            end_run = i
            end_off = pos + len(token) - cum
            break
        cum += rlen

    if start_run is None or end_run is None:
        return False

    sr = runs[start_run]
    er = runs[end_run]
    before = sr.text[:start_off]
    after = er.text[end_off:]

    sr.text = before + first_line
    if start_run != end_run:
        er.text = after
    else:
        sr.text = before + first_line + after

    for i in range(start_run + 1, end_run + (1 if start_run != end_run else 0)):
        pass
    if start_run != end_run:
        for i in range(end_run - 1, start_run, -1):
            r_el = runs[i]._r
            r_el.getparent().remove(r_el)

    if rest_lines:
        template_p_el = paragraph._p
        parent = template_p_el.getparent()
        insert_at = list(parent).index(template_p_el) + 1
        template_rpr = sr._r.find(qn("a:rPr"))
        for line in rest_lines:
            new_p_el = copy.deepcopy(template_p_el)
            for r_el in new_p_el.findall(qn("a:r")):
                new_p_el.remove(r_el)
            new_r_el = new_p_el.makeelement(qn("a:r"), {})
            if template_rpr is not None:
                new_r_el.append(copy.deepcopy(template_rpr))
            t_el = new_r_el.makeelement(qn("a:t"), {})
            t_el.text = line
            new_r_el.append(t_el)
            new_p_el.append(new_r_el)
            parent.insert(insert_at, new_p_el)
            insert_at += 1

    return True


def ensure_paragraph_spacing(shape, token):
    """Add a small breathing-room space-after on the placeholder's paragraph
    so that when the category/nominee text wraps to multiple lines, it does
    not visually collide with the box below it (e.g. Award Category wrapping
    into the Nominees list). Only applied to the category/zone placeholders,
    since the nominee/winner list already gets its own per-line spacing."""
    if token not in ("<<award category>>", "<<zone>>"):
        return
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        pPr = paragraph._p.find(qn("a:pPr"))
        if pPr is None:
            pPr = paragraph._p.makeelement(qn("a:pPr"), {})
            paragraph._p.insert(0, pPr)
        for old_tag in ("a:spcAft",):
            existing = pPr.find(qn(old_tag))
            if existing is not None:
                pPr.remove(existing)
        spcAft = pPr.makeelement(qn("a:spcAft"), {})
        spcPts = spcAft.makeelement(qn("a:spcPts"), {"val": "1200"})
        spcAft.append(spcPts)
        pPr.append(spcAft)


def fill_shapes_recursive(shapes, mapping):
    for shape in shapes:
        if shape.shape_type == 6 and hasattr(shape, "shapes"):
            fill_shapes_recursive(shape.shapes, mapping)
            continue
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue
        for paragraph in list(shape.text_frame.paragraphs):
            low_text = paragraph.text.lower()
            for token, key in TOKENS.items():
                if token in low_text:
                    replace_token_in_run_list(paragraph, token, mapping.get(key, ""))
                    ensure_paragraph_spacing(shape, token)
                    break


def fill_slide(slide, mapping):
    fill_shapes_recursive(slide.shapes, mapping)


def build_deck(excel_path, template_path):
    rows = read_rows(excel_path)
    nominee_groups, winner_groups, key_order = group_rows(rows)

    prs = Presentation(template_path)
    nominee_idx, winner_idx = pick_template_indices(prs)

    # Preserve the exact order categories first appear in the Excel file
    # (this drives the emcee's announcement sequence), rather than sorting
    # alphabetically. Within each category, the nominee slide (if any)
    # comes first, immediately followed by that category's winner slide(s).
    keys = key_order

    for key in keys:
        zone, category = key
        if key in nominee_groups:
            entries = nominee_groups[key]
            new_slide = duplicate_slide(prs, nominee_idx)
            fill_slide(new_slide, {
                "ZONE": zone,
                "AWARD CATEGORY": category,
                "NOMINEES": "\n".join(e["nominee_name"] for e in entries),
                "NOMINEES_WORD": "NOMINEES",
                "PLACEHOLDER_X": entries[0]["placeholder_x"],
                "PLACEHOLDER_Y": entries[0]["placeholder_y"],
                "PLACEHOLDER_Z": entries[0]["placeholder_z"],
            })
        if key in winner_groups:
            entries = winner_groups[key]
            # Winners always get one slide per winner, even when multiple
            # winners share the same Zone + Award Category (e.g. joint
            # winners or multiple sub-winners), so each gets full visibility.
            for entry in entries:
                new_slide = duplicate_slide(prs, winner_idx)
                fill_slide(new_slide, {
                    "ZONE": zone,
                    "AWARD CATEGORY": category,
                    "WINNER": entry["winner_name"],
                    "WINNER_WORD": "WINNER",
                    "PLACEHOLDER_X": entry["placeholder_x"],
                    "PLACEHOLDER_Y": entry["placeholder_y"],
                    "PLACEHOLDER_Z": entry["placeholder_z"],
                })

    for idx in sorted({nominee_idx, winner_idx}, reverse=True):
        delete_slide(prs, idx)

    return prs


def main():
    if len(sys.argv) != 4:
        print("Usage: python generate_award_slides.py input.xlsx template.pptx output.pptx")
        raise SystemExit(1)
    deck = build_deck(sys.argv[1], sys.argv[2])
    deck.save(sys.argv[3])
    print(f"Saved {sys.argv[3]} ({len(deck.slides)} slides)")


if __name__ == "__main__":
    main()
