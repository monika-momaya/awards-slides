from pathlib import Path
import shutil, zipfile
out = Path('/home/user/output/streamlit_award_app_v4')
if out.exists(): shutil.rmtree(out)
out.mkdir(parents=True, exist_ok=True)

app = r'''import streamlit as st
from pathlib import Path
from copy import deepcopy
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
import re
import csv
from io import BytesIO

st.set_page_config(page_title='Award Slide Generator', layout='wide')

PLACEHOLDERS = {
    'winner_word': '<<winner-word>>',
    'nominees_word': '<<nominees-word>>',
    'award_category': '<<award category>>',
    'zone': '<<zone>>',
    'nominees': '<<nominees>>',
    'winner': '<<winner>>',
}
ZONE_WORDS = ['south', 'west', 'north', 'east', 'pan india']


def s(v):
    return '' if v is None else str(v).strip()


def norm(v):
    return re.sub(r'\s+', ' ', s(v)).strip().lower()


def load_rows(uploaded_file):
    wb = load_workbook(uploaded_file, data_only=True)
    ws = wb[wb.sheetnames[0]]
    return list(ws.iter_rows(values_only=True))


def detect_style(rows):
    header = [norm(x) for x in rows[0]] if rows else []
    joined = ' | '.join(header)
    if 'nominees' in joined and 'award category' in joined:
        return 'nominee'
    if 'winner' in joined or 'winner category' in joined or 'bo name' in joined:
        return 'winner'
    return 'unknown'


def extract_categories(rows):
    header_row = 0
    headers = [norm(x) for x in rows[header_row]] if rows else []
    idx_cat = next((i for i,h in enumerate(headers) if 'award category' in h or 'winner category' in h or h == 'category'), 0)
    idx_name = next((i for i,h in enumerate(headers) if 'nominees' in h or 'bo name' in h or 'name' in h), 1)
    idx_zone = next((i for i,h in enumerate(headers) if 'zone' in h), None)
    recs = []
    current = None
    for row in rows[header_row+1:]:
        vals = [s(x) for x in row]
        if not any(vals):
            continue
        cat = vals[idx_cat] if idx_cat < len(vals) else ''
        name = vals[idx_name] if idx_name < len(vals) else ''
        zone = vals[idx_zone] if idx_zone is not None and idx_zone < len(vals) else ''
        if cat:
            current = {'category': cat, 'zone': zone, 'nominees': [], 'winner': ''}
            recs.append(current)
        if current and name:
            current['nominees'].append(name)
    return recs


def extract_winners(rows):
    header = [norm(x) for x in rows[0]] if rows else []
    idx_cat = next((i for i,h in enumerate(header) if 'winner category' in h or 'award category' in h or h == 'category'), 1)
    idx_winner = next((i for i,h in enumerate(header) if 'bo name' in h or 'winner' in h or 'name' in h), 0)
    idx_zone = next((i for i,h in enumerate(header) if 'zone' in h), None)
    recs = []
    for row in rows[1:]:
        vals = [s(x) for x in row]
        if not any(vals):
            continue
        cat = vals[idx_cat] if idx_cat < len(vals) else ''
        winner = vals[idx_winner] if idx_winner < len(vals) else ''
        zone = vals[idx_zone] if idx_zone is not None and idx_zone < len(vals) else ''
        if cat and winner:
            recs.append({'category': cat, 'zone': zone, 'nominees': [], 'winner': winner})
    return recs


def parse_source(uploaded_file):
    rows = load_rows(uploaded_file)
    style = detect_style(rows)
    if style == 'nominee':
        return extract_categories(rows)
    if style == 'winner':
        return extract_winners(rows)
    # fallback: try both
    recs = extract_categories(rows)
    if recs:
        return recs
    return extract_winners(rows)


def set_text(shape, text):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = s(text)
    r.font.size = Pt(22)


def fill_slide(slide, rec):
    mapping = {
        PLACEHOLDERS['winner_word']: 'Winner',
        PLACEHOLDERS['nominees_word']: 'Nominees',
        PLACEHOLDERS['award_category']: rec.get('category', ''),
        PLACEHOLDERS['zone']: rec.get('zone', ''),
        PLACEHOLDERS['nominees']: '\n'.join(rec.get('nominees', [])),
        PLACEHOLDERS['winner']: rec.get('winner', ''),
    }
    for sh in slide.shapes:
        if not hasattr(sh, 'text'):
            continue
        txt = norm(sh.text)
        for ph, val in mapping.items():
            if ph in txt:
                set_text(sh, val)


def build_ppt(template_bytes, recs):
    prs = Presentation(BytesIO(template_bytes))
    out = Presentation()
    out.slide_width = prs.slide_width
    out.slide_height = prs.slide_height
    for sldId in list(out.slides._sldIdLst):
        out.slides._sldIdLst.remove(sldId)
    if not prs.slides:
        raise RuntimeError('Template has no slides')
    template_slide = prs.slides[0]
    for rec in recs:
        new_slide = out.slides.add_slide(out.slide_layouts[6])
        for shape in list(new_slide.shapes):
            el = shape.element
            el.getparent().remove(el)
        for shape in template_slide.shapes:
            newel = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        fill_slide(new_slide, rec)
    return out


def recs_to_csv_bytes(recs):
    import io
    txt = io.StringIO()
    w = csv.DictWriter(txt, fieldnames=['category', 'zone', 'winner', 'nominees'])
    w.writeheader()
    for r in recs:
        w.writerow({'category': r.get('category', ''), 'zone': r.get('zone', ''), 'winner': r.get('winner', ''), 'nominees': ' | '.join(r.get('nominees', []))})
    return txt.getvalue().encode('utf-8')

st.title('Award Slide Generator')
st.write('Two uploads only: Excel and PowerPoint template.')
source = st.file_uploader('Upload Excel file', type=['xlsx'])
template = st.file_uploader('Upload PowerPoint template', type=['pptx'])

if source and template and st.button('Generate'):
    try:
        recs = parse_source(source)
        if not recs:
            st.error('No categories or winners could be detected from the Excel file.')
        else:
            prs = build_ppt(template.getvalue(), recs)
            bio = BytesIO()
            prs.save(bio)
            st.success(f'Generated {len(recs)} slide(s).')
            st.download_button('Download PPTX', bio.getvalue(), file_name='OUTPUT_awards.pptx', mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            st.download_button('Download CSV', recs_to_csv_bytes(recs), file_name='award_records.csv', mime='text/csv')
    except Exception as e:
        st.error(str(e))
'''

readme = '''Streamlit Award Slide Generator v4\n\nRule:\n- one slide per category\n- nominees appear one below the other on the slide\n- one slide per winner row when using winner source\n\nUse only two uploads: Excel and PowerPoint template.\n'''
req = 'streamlit\npython-pptx\nopenpyxl\n'
(out/'award_streamlit_app.py').write_text(app, encoding='utf-8')
(out/'README.txt').write_text(readme, encoding='utf-8')
(out/'requirements.txt').write_text(req, encoding='utf-8')
shutil.copy('/home/user/Template-2.pptx', out/'Template-2.pptx')
zip_path = Path('/home/user/output/streamlit_award_app_v4.zip')
if zip_path.exists(): zip_path.unlink()
with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
    for p in out.iterdir():
        if p.is_file(): zf.write(p, arcname=p.name)
print(str(zip_path))